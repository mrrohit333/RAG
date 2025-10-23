import os
import fitz  # PyMuPDF
import pandas as pd
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
import ollama
from langchain.text_splitter import RecursiveCharacterTextSplitter
import torch
import pickle
import traceback
from datetime import datetime
import json

# ==== CONFIG ====
BASE_DIR = "data"
DIM = 384
MODEL_NAME = "all-MiniLM-L6-v2"
RELEVANCE_THRESHOLD = 1.2
METADATA_FILE = "metadata.json"

# ==== INIT ====
device = "cuda" if torch.cuda.is_available() else "cpu"
print(f"✅ Using device: {device}")
embed_model = SentenceTransformer(MODEL_NAME, device=device)


# -------------------------------------------------
# Utility Functions
# -------------------------------------------------
def get_user_dirs(user_id):
    """Return user directory, chunk cache path, and FAISS index path."""
    user_dir = os.path.join(BASE_DIR, str(user_id))
    vector_cache = os.path.join(user_dir, "chunks.pkl")
    index_path = os.path.join(user_dir, "faiss.index")
    os.makedirs(user_dir, exist_ok=True)
    return user_dir, vector_cache, index_path


def update_metadata(user_id, file_path, num_chunks):
    """Store metadata for uploaded documents (filename, time, chunk count)."""
    user_dir, _, _ = get_user_dirs(user_id)
    metadata_path = os.path.join(user_dir, METADATA_FILE)

    metadata = []
    if os.path.exists(metadata_path):
        try:
            with open(metadata_path, "r", encoding="utf-8") as f:
                metadata = json.load(f)
        except Exception:
            metadata = []

    metadata.append({
        "file": os.path.basename(file_path),
        "chunks": num_chunks,
        "uploaded_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    })

    with open(metadata_path, "w", encoding="utf-8") as f:
        json.dump(metadata, f, indent=2)
    print(f"🗂️ Metadata updated for user {user_id} ({len(metadata)} total uploads)")


def chunk_file(file_path):
    """Extract text from supported files and split into semantic chunks."""
    try:
        filename = os.path.basename(file_path)
        ext = filename.lower().split(".")[-1]
        text = ""

        if ext == "pdf":
            print(f"📄 Loading PDF: {filename}")
            doc = fitz.open(file_path)
            text = "\n".join(page.get_text() for page in doc)
        elif ext == "txt":
            print(f"📜 Loading TXT: {filename}")
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        elif ext == "csv":
            print(f"📊 Loading CSV: {filename}")
            df = pd.read_csv(file_path)
            text = "\n".join(df.astype(str).apply(lambda x: " | ".join(x), axis=1))
        elif ext == "tsv":
            print(f"📊 Loading TSV: {filename}")
            df = pd.read_csv(file_path, sep="\t")
            text = "\n".join(df.astype(str).apply(lambda x: " | ".join(x), axis=1))
        elif ext == "xlsx":
            print(f"📊 Loading XLSX: {filename}")
            df = pd.read_excel(file_path, engine="openpyxl")
            text = "\n".join(df.astype(str).apply(lambda x: " | ".join(x), axis=1))
        elif ext == "docx":
            print(f"📝 Loading DOCX: {filename}")
            try:
                from docx import Document
                doc = Document(file_path)
                text = "\n".join(p.text for p in doc.paragraphs)
            except Exception as ex:
                print(f"❌ Failed to parse DOCX: {ex}")
                return []
        elif ext == "pptx":
            print(f"🖼️ Loading PPTX: {filename}")
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                slides_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slides_text.append(shape.text)
                text = "\n".join(slides_text)
            except Exception as ex:
                print(f"❌ Failed to parse PPTX: {ex}")
                return []
        else:
            print(f"⚠️ Unsupported file type: {ext}")
            return []

        if not text.strip():
            print(f"⚠️ No readable text found in {filename}")
            return []

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=800,
            chunk_overlap=100,
            separators=["\n\n", "\n", ".", "!", "?", " "]
        )
        chunks = splitter.split_text(text)
        print(f"✅ {len(chunks)} chunks created from {filename}")
        return chunks

    except Exception as e:
        print(f"❌ Error while chunking {file_path}: {str(e)}")
        traceback.print_exc()
        return []


def build_faiss_index(chunks):
    """Build FAISS index from chunks."""
    try:
        if not chunks:
            return None, None, None
        print(f"🔢 Encoding {len(chunks)} chunks into embeddings...")
        embeddings = embed_model.encode(chunks, convert_to_numpy=True, show_progress_bar=False)
        index = faiss.IndexFlatL2(DIM)
        index.add(np.array(embeddings))
        return index, embeddings, chunks
    except Exception as e:
        print(f"❌ Error building FAISS index: {str(e)}")
        traceback.print_exc()
        return None, None, None


def save_user_index(user_id, index, chunks):
    """Save FAISS index and chunks for a user."""
    try:
        user_dir, vector_cache, index_path = get_user_dirs(user_id)
        with open(vector_cache, "wb") as f:
            pickle.dump(chunks, f)
        faiss.write_index(index, index_path)
        print(f"💾 Saved FAISS index & {len(chunks)} chunks for user {user_id}")
    except Exception as e:
        print(f"❌ Failed to save index: {str(e)}")
        traceback.print_exc()


def clear_user_index(user_id):
    """Remove index and chunk cache files for a user if present."""
    user_dir, vector_cache, index_path = get_user_dirs(user_id)
    try:
        if os.path.exists(vector_cache):
            os.remove(vector_cache)
        if os.path.exists(index_path):
            os.remove(index_path)
        print(f"🧹 Cleared vector cache and index for user {user_id}")
    except Exception as e:
        print(f"⚠️ Failed to clear index files: {str(e)}")


def rebuild_index_from_files(user_id):
    """
    Rebuild user's FAISS index by re-processing all remaining files listed in metadata.json.
    If no files remain, clears the index/chunk cache.
    """
    user_dir, _, _ = get_user_dirs(user_id)
    metadata_path = os.path.join(user_dir, METADATA_FILE)
    if not os.path.exists(metadata_path):
        clear_user_index(user_id)
        return True

    try:
        with open(metadata_path, "r", encoding="utf-8") as f:
            metadata = json.load(f)
    except Exception as e:
        print(f"⚠️ Failed to read metadata for rebuild: {str(e)}")
        metadata = []

    all_chunks = []
    new_metadata = []
    for entry in metadata:
        file_path = os.path.join(user_dir, entry.get("file", ""))
        if not os.path.exists(file_path):
            continue
        chunks = chunk_file(file_path)
        if chunks:
            all_chunks.extend(chunks)
            new_metadata.append({
                "file": os.path.basename(file_path),
                "chunks": len(chunks),
                "uploaded_at": entry.get("uploaded_at") or datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            })

    # Update metadata to reflect current chunk counts and remove missing files
    with open(metadata_path, "w", encoding="utf-8") as f:
        json.dump(new_metadata, f, indent=2)

    if not all_chunks:
        clear_user_index(user_id)
        return True

    index, _, _ = build_faiss_index(all_chunks)
    if index is None:
        print("❌ Rebuild failed: no index created.")
        return False
    save_user_index(user_id, index, all_chunks)
    print(f"✅ Rebuilt index for user {user_id} from {len(new_metadata)} documents, {len(all_chunks)} chunks")
    return True


def delete_user_document(user_id: int, filename: str) -> bool:
    """
    Delete a specific document for a user and rebuild the FAISS index from remaining documents.
    Safe even when metadata is missing or file is absent.
    """
    user_dir, _, _ = get_user_dirs(user_id)
    target_path = os.path.join(user_dir, filename)

    # Remove file if exists
    try:
        if os.path.exists(target_path):
            os.remove(target_path)
            print(f"🗑️ Deleted file: {target_path}")
    except Exception as e:
        print(f"⚠️ Could not delete file {target_path}: {str(e)}")

    # Update metadata to remove this entry
    metadata_path = os.path.join(user_dir, METADATA_FILE)
    if os.path.exists(metadata_path):
        try:
            with open(metadata_path, "r", encoding="utf-8") as f:
                metadata = json.load(f)
            metadata = [m for m in metadata if m.get("file") != filename]
            with open(metadata_path, "w", encoding="utf-8") as f:
                json.dump(metadata, f, indent=2)
        except Exception as e:
            print(f"⚠️ Failed to update metadata after delete: {str(e)}")

    # Rebuild index and chunks from remaining files
    return rebuild_index_from_files(user_id)
def load_user_index(user_id):
    """Load FAISS index and chunks for a user."""
    user_dir, vector_cache, index_path = get_user_dirs(user_id)
    if not os.path.exists(index_path) or not os.path.exists(vector_cache):
        return None, None
    try:
        index = faiss.read_index(index_path)
        with open(vector_cache, "rb") as f:
            chunks = pickle.load(f)
        return index, chunks
    except Exception as e:
        print(f"⚠️ Failed to load user index: {str(e)} — will rebuild.")
        traceback.print_exc()
        return None, None


# -------------------------------------------------
# Document Upload & Index Building (Multiple Docs)
# -------------------------------------------------
def add_new_document(file_path, user_id):
    """Add new document chunks to user's FAISS index (append-safe)."""
    try:
        print(f"📂 Starting document processing for user {user_id}: {file_path}")
        chunks = chunk_file(file_path)
        if not chunks:
            print("⚠️ No chunks generated — skipping.")
            return False

        user_dir, vector_cache, index_path = get_user_dirs(user_id)

        # Load existing data if available
        existing_chunks = []
        index = None
        if os.path.exists(index_path):
            index, existing_chunks = load_user_index(user_id)
            if index is None:
                print("⚠️ Corrupted index detected. Rebuilding clean index...")
                all_chunks = (existing_chunks or []) + chunks
                index, _, _ = build_faiss_index(all_chunks)
                if index:
                    save_user_index(user_id, index, all_chunks)
                    update_metadata(user_id, file_path, len(chunks))
                    print("✅ Rebuilt full FAISS index.")
                    return True
                else:
                    print("❌ Failed to rebuild index.")
                    return False
        else:
            index, _, _ = build_faiss_index(chunks)

        # Create embeddings for new chunks and append
        if index is not None:
            print(f"🔢 Adding {len(chunks)} new chunks to FAISS index...")
            new_embeddings = embed_model.encode(chunks, convert_to_numpy=True, show_progress_bar=False)
            index.add(np.array(new_embeddings))
            all_chunks = (existing_chunks or []) + chunks
            save_user_index(user_id, index, all_chunks)
            update_metadata(user_id, file_path, len(chunks))
            print(f"✅ Successfully added {len(chunks)} chunks from {os.path.basename(file_path)}")
            return True
        else:
            print("❌ Could not create FAISS index.")
            return False

    except Exception as e:
        print(f"❌ add_new_document() failed: {str(e)}")
        traceback.print_exc()
        return False


# -------------------------------------------------
# Query Handling (Safe + Auto Repair)
# -------------------------------------------------
def get_response_stream(user_id, query):
    """Generate streamed answers from user's documents (auto-repair + safe)."""
    print(f"💬 Query from user {user_id}: {query}")

    index, chunks = load_user_index(user_id)

    # 🧱 Auto-repair if vector/chunk mismatch
    if index is not None:
        total_vectors = index.ntotal
        chunk_count = len(chunks) if chunks else 0
        if total_vectors != chunk_count:
            print(f"⚠️ Detected mismatch: {total_vectors} vectors vs {chunk_count} chunks.")
            print("🧱 Rebuilding FAISS index automatically...")
            if chunks:
                index, _, _ = build_faiss_index(chunks)
                if index:
                    save_user_index(user_id, index, chunks)
                    print("✅ Index rebuilt successfully.")
                else:
                    print("❌ Index rebuild failed.")
                    index = None

    # 🧠 Construct query prompt
    if index is None or not chunks:
        prompt = f"You have no uploaded documents. Answer this generally:\n\nQ: {query}\nA:"
    else:
        query_vec = embed_model.encode([query], convert_to_numpy=True)
        D, I = index.search(np.array(query_vec), k=3)
        valid_indices = [i for i in I[0] if i < len(chunks)]

        if not valid_indices or D[0][0] > RELEVANCE_THRESHOLD:
            prompt = f"No relevant info found in your uploaded docs. Answer generally:\n\nQ: {query}\nA:"
        else:
            context = "\n\n".join([chunks[i] for i in valid_indices])
            prompt = f"Use the context below to answer:\n\nContext:\n{context}\n\nQ: {query}\nA:"

    # 🔄 Stream answer from Ollama model
    stream = ollama.chat(
        model="llama3:8b",
        messages=[{"role": "user", "content": prompt}],
        stream=True
    )

    for chunk in stream:
        if "message" in chunk and "content" in chunk["message"]:
            yield chunk["message"]["content"]
